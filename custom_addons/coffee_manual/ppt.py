from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create a NEW presentation (not the old one)
prs = Presentation()

# Define a simple layout for all slides
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]

# --- Monthly Report Content ---
slides_content = [
    {
        "title": "Monthly Report (October 15 – November 5, 2025)",
        "subtitle": "ICT: System Development Division\nAMG HOLDINGS"
    },
    {
        "title": "CRM System (Fully Completed)",
        "content": (
            "The CRM system has reached full operational status. The recent phase focused on:\n\n"
            "- Integrating custom CRM modules into a unified marketing system.\n"
            "- Embedding data quality and telemarketing confirmations into core lead workflows.\n"
            "- Automating KPI computation and scoring through real-time updates.\n"
            "- Developing role-based access control for sales, marketing, and management users.\n"
            "- Conducting comprehensive testing, presentation, and feedback sessions with the ICT department.\n\n"
            "Next Steps:\n"
            "- Address ICT feedback on the CRM Business Requirement Document.\n"
            "- Initiate user-level testing with actual business data prior to deployment.\n"
            "- Prepare rollout documentation and deployment roadmap."
        )
    },
    {
        "title": "KPI Framework (Completed and Integrated)",
        "content": (
            "The KPI Framework was finalized and merged with the CRM system to provide:\n\n"
            "- Dynamic KPI definition and tracking for data quality and lead management.\n"
            "- Automatic linkage between telemarketing activities and KPI achievements.\n"
            "- Real-time computation for both actual and target values, improving performance analysis.\n\n"
            "Outcome:\n"
            "Improved decision-making visibility and precise evaluation of individual and departmental productivity."
        )
    },
    {
        "title": "Active Directory (AD) System Recovery",
        "content": (
            "A system-wide domain user lockout occurred, impacting all users. The issue was thoroughly investigated "
            "and resolved after extensive recovery efforts over two days. All domain users have since regained normal operations.\n\n"
            "Actions Taken:\n"
            "- Restored all domain accounts and validated group policies.\n"
            "- Introduced monitoring and access control measures to prevent recurrence.\n\n"
            "Outcome:\n"
            "Full user access restored and AD system stabilized."
        )
    },
    {
        "title": "Other Activities",
        "content": (
            "- Inventory Module: Initiated product master data setup with ICT collaboration.\n"
            "- System Testing: Verified integration across CRM, KPI, telemarketing, competitor log, and marketing intelligence modules.\n"
            "- Documentation: Prepared detailed system documentation and reports for upcoming presentation.\n"
            "- Weekly & Monthly Meetings: Participated in ICT coordination and planning sessions."
        )
    },
    {
        "title": "Challenges and Resolutions",
        "content": (
            "Challenges:\n"
            "- User Lockout Incident: Required deep technical troubleshooting and multi-day recovery efforts.\n"
            "- Module Consolidation: Merging standalone CRM modules required refinement to maintain logic consistency.\n\n"
            "Resolutions:\n"
            "- Implemented stricter credential policies and AD monitoring tools.\n"
            "- Standardized CRM logic to prevent computation inconsistencies."
        )
    },
    {
        "title": "Next Month’s Focus",
        "content": (
            "- Conduct final user acceptance testing (UAT) for CRM and KPI modules.\n"
            "- Deploy the unified Marketing System.\n"
            "- Initiate training and documentation handover to end-user departments.\n"
            "- Begin development of the JS-based price comparison dashboard for AMG vs. competitor data visualization."
        )
    },
]

# --- Helper Function to Add Slides ---
def add_slide(title, content=None, subtitle=None):
    layout = title_slide_layout if subtitle and not content else content_slide_layout
    slide = prs.slides.add_slide(layout)

    # Title
    if slide.shapes.title:
        slide.shapes.title.text = title

    # Subtitle (for title slide)
    if subtitle:
        placeholders = [ph for ph in slide.placeholders if ph.placeholder_format.idx != 0]
        if placeholders:
            placeholders[0].text = subtitle
        else:
            left, top, width, height = Inches(1), Inches(2), Inches(8), Inches(1)
            box = slide.shapes.add_textbox(left, top, width, height)
            tf = box.text_frame
            p = tf.add_paragraph()
            p.text = subtitle
            p.font.size = Pt(18)
            p.alignment = PP_ALIGN.CENTER

    # Body (content)
    if content:
        left, top, width, height = Inches(1), Inches(1.5), Inches(8.5), Inches(5)
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = content
        p.font.size = Pt(16)
        p.alignment = PP_ALIGN.LEFT


# --- Generate Slides ---
for s in slides_content:
    add_slide(s["title"], s.get("content"), s.get("subtitle"))

# Save File
output_path = "D:/crm_project/custom_addons/coffee_manual/Monthly_Report_Oct15_Nov5_2025.pptx"
prs.save(output_path)

print(f"✅ Monthly report generated successfully:\n{output_path}")
